from __future__ import annotations

import json
import shutil
from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = REPO_ROOT / "reports"
FIG_DIR = REPORTS_DIR / "presentation_figures_v2"
FIG_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = REPO_ROOT / "Interim_Presentation_Draft.pptx"

GENERATED_ASSET = Path(
    "/Users/gregoiredambre/.codex/generated_images/019dfd23-be51-7d81-984b-f1fb3315788d/"
    "ig_076a8ad7da1411a4016a158034bb90819183bdce9db19b8e13.png"
)
LOCAL_ASSET = FIG_DIR / "from_asset_to_network.png"

CAMBRIDGE_BLUE = "#003B6F"
CAMBRIDGE_BLUE_LIGHT = "#6A8FB3"
SLATE = "#48566A"
TEAL = "#5C8D89"
LIGHT_GREY = "#D9DFE5"
TEXT = "#1F2933"
RED = "#A33A3A"
GREEN = "#2E7D32"
ORANGE = "#C27A2C"


def hex_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def setup_matplotlib() -> None:
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 11,
            "axes.titlesize": 14,
            "axes.labelsize": 11,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "figure.facecolor": "white",
            "axes.facecolor": "white",
            "savefig.facecolor": "white",
            "savefig.dpi": 220,
            "savefig.bbox": "tight",
        }
    )


def copy_problem_statement_visual() -> Path | None:
    if GENERATED_ASSET.exists():
        shutil.copy2(GENERATED_ASSET, LOCAL_ASSET)
        return LOCAL_ASSET
    return None


def build_hpms16_benchmark() -> Path:
    ensemble = json.loads((REPO_ROOT / "graph_data" / "ensemble_results.json").read_text())
    hpms_materials = json.loads((REPORTS_DIR / "materials_weight_sweep_hpms16.json").read_text())

    rows = [
        ("Random Forest", ensemble["results"]["RF local"]["test"]["r2"], CAMBRIDGE_BLUE_LIGHT),
        ("R-GCN\n(full refined)", ensemble["results"]["R-GCN"]["test"]["r2"], SLATE),
        ("R-GCN refined\n(materials + weights)", hpms_materials["best_test_r2"], TEAL),
        ("Stacked MLP\nensemble", ensemble["results"]["Stacked MLP"]["test"]["r2"], CAMBRIDGE_BLUE),
    ]
    df = pd.DataFrame(rows, columns=["model", "r2", "color"]).sort_values("r2")

    fig, ax = plt.subplots(figsize=(9.8, 5.2))
    bars = ax.barh(df["model"], df["r2"], color=df["color"], edgecolor="white", height=0.68)
    for bar, value in zip(bars, df["r2"]):
        ax.text(value + 0.006, bar.get_y() + bar.get_height() / 2, f"{value:.3f}", va="center", fontweight="bold")
    ax.set_xlim(0, 0.62)
    ax.set_xlabel("Test R²")
    ax.set_title("HPMS16 cracking: benchmark on the main comparative target", loc="left", color=CAMBRIDGE_BLUE)
    ax.grid(axis="x", color=LIGHT_GREY, linewidth=0.7)
    ax.set_axisbelow(True)

    out = FIG_DIR / "hpms16_benchmark.png"
    fig.savefig(out)
    plt.close(fig)
    return out


def build_mepdg_benchmark() -> Path:
    df = pd.read_csv(REPORTS_DIR / "mepdg_benchmark.csv")
    colors = {
        "RF local": CAMBRIDGE_BLUE_LIGHT,
        "R-GCN baseline": SLATE,
        "Stacked MLP ensemble": CAMBRIDGE_BLUE,
        "R-GCN best materials sweep (climate_pavement)": TEAL,
    }
    labels = {
        "RF local": "Random Forest",
        "R-GCN baseline": "R-GCN\nbaseline",
        "Stacked MLP ensemble": "Stacked MLP\nensemble",
        "R-GCN best materials sweep (climate_pavement)": "R-GCN refined\n(materials + weights)",
    }
    df["label"] = df["model"].map(labels)
    df["color"] = df["model"].map(colors)
    df = df.sort_values("r2_test")

    fig, ax = plt.subplots(figsize=(9.8, 5.2))
    bars = ax.barh(df["label"], df["r2_test"], color=df["color"], edgecolor="white", height=0.68)
    for bar, value in zip(bars, df["r2_test"]):
        ax.text(value + 0.006, bar.get_y() + bar.get_height() / 2, f"{value:.3f}", va="center", fontweight="bold")
    ax.set_xlim(0, 0.64)
    ax.set_xlabel("Test R²")
    ax.set_title("MEPDG cracking: refined graph model gives the best absolute score", loc="left", color=CAMBRIDGE_BLUE)
    ax.grid(axis="x", color=LIGHT_GREY, linewidth=0.7)
    ax.set_axisbelow(True)

    out = FIG_DIR / "mepdg_benchmark.png"
    fig.savefig(out)
    plt.close(fig)
    return out


def build_ood_figure() -> Path:
    df = pd.read_csv(REPORTS_DIR / "part1_ood_ensemble.csv")
    states = df["state_held_out"].astype(str).tolist()
    x = np.arange(len(states))
    width = 0.24

    fig, ax = plt.subplots(figsize=(10.2, 5.5))
    ax.bar(x - width, df["rf_test_r2"], width, label="Random Forest", color=CAMBRIDGE_BLUE_LIGHT)
    ax.bar(x, df["rgcn_test_r2"], width, label="R-GCN", color=SLATE)
    ax.bar(x + width, df["ensemble_test_r2"], width, label="Stacked MLP ensemble", color=CAMBRIDGE_BLUE)
    ax.axhline(0, color=TEXT, linewidth=0.8)
    ax.set_xticks(x)
    ax.set_xticklabels(states)
    ax.set_xlabel("Held-out state")
    ax.set_ylabel("Test R²")
    ax.set_title("Geographic transfer remains weak out-of-domain", loc="left", color=CAMBRIDGE_BLUE)
    ax.legend(loc="upper right")
    ax.grid(axis="y", color=LIGHT_GREY, linewidth=0.7)
    ax.set_axisbelow(True)

    summary = json.loads((REPORTS_DIR / "part1_ood_ensemble_summary.json").read_text())
    caption = (
        f"Mean R² over {summary['n_states_evaluated']} representative states: "
        f"RF {summary['rf_r2']['mean']:.3f}, ensemble {summary['ensemble_r2']['mean']:.3f}."
    )
    fig.text(0.01, -0.03, caption, fontsize=10, color=SLATE)

    out = FIG_DIR / "ood_ensemble.png"
    fig.savefig(out)
    plt.close(fig)
    return out


def build_traffic_figure() -> Path:
    payload = json.loads((REPORTS_DIR / "traffic_sensitivity_bootstrap.json").read_text())
    labels = ["Maintenance\nevents", "Controls"]
    means = [payload["event_mean"], payload["control_mean"]]
    low = [payload["event_mean"] - payload["event_ci_low"], payload["control_mean"] - payload["control_ci_low"]]
    high = [payload["event_ci_high"] - payload["event_mean"], payload["control_ci_high"] - payload["control_mean"]]
    yerr = np.array([low, high])

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.8, 5.2), gridspec_kw={"width_ratios": [1.3, 1.0]})

    bars = ax1.bar(labels, means, yerr=yerr, capsize=6, color=[CAMBRIDGE_BLUE, SLATE], edgecolor="white")
    for bar, value in zip(bars, means):
        ax1.text(bar.get_x() + bar.get_width() / 2, value + 500, f"{value:,.0f}", ha="center", va="bottom", fontweight="bold")
    ax1.set_ylabel("Mean pre→post delta of ANNUAL_GESAL_TREND")
    ax1.set_title("Maintenance years leave a measurable annual traffic-loading signature", loc="left", color=CAMBRIDGE_BLUE)
    ax1.grid(axis="y", color=LIGHT_GREY, linewidth=0.7)
    ax1.set_axisbelow(True)

    diff = payload["diff_mean"]
    ci_low = payload["diff_ci_low"]
    ci_high = payload["diff_ci_high"]
    ax2.errorbar([0], [diff], yerr=[[diff - ci_low], [ci_high - diff]], fmt="o", color=RED, markersize=10, linewidth=2.2, capsize=8)
    ax2.axhline(0, color=TEXT, linewidth=0.8)
    ax2.set_xlim(-0.8, 0.8)
    ax2.set_xticks([0])
    ax2.set_xticklabels(["Events - controls"])
    ax2.set_ylabel("Difference in mean delta")
    ax2.set_title(f"Bootstrap IC95% excludes zero\n(p = {payload['welch_p_value']:.4f})", color=CAMBRIDGE_BLUE)
    ax2.grid(axis="y", color=LIGHT_GREY, linewidth=0.7)
    ax2.set_axisbelow(True)

    out = FIG_DIR / "traffic_bootstrap.png"
    fig.savefig(out)
    plt.close(fig)
    return out


def add_textbox(slide, left, top, width, height, text, font_size=20, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_rgb(color)
    p.alignment = align
    return box


def add_bullet_box(slide, left, top, width, height, bullets, font_size=20, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_rgb(color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.5), title, font_size=28, bold=True, color=CAMBRIDGE_BLUE)
    if subtitle:
        add_textbox(slide, Inches(0.62), Inches(0.8), Inches(11.4), Inches(0.35), subtitle, font_size=12, color=SLATE)


def create_presentation(figs: dict[str, Path], problem_image: Path | None) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_rgb("FFFFFF")
    add_textbox(slide, Inches(0.7), Inches(0.9), Inches(10.5), Inches(0.8), "Interim Presentation", font_size=30, bold=True, color=CAMBRIDGE_BLUE)
    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(11), Inches(0.6), "Graph neural networks for pavement deterioration prediction and network-aware maintenance planning", font_size=22, color=TEXT)
    add_textbox(slide, Inches(0.72), Inches(2.6), Inches(8), Inches(0.4), "MPhil ISMM Cambridge", font_size=16, color=SLATE)
    add_textbox(slide, Inches(0.72), Inches(2.95), Inches(8), Inches(0.4), "Gregoire Dambre", font_size=16, color=SLATE)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(5.9), Inches(2.3), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = hex_rgb(CAMBRIDGE_BLUE)
    accent.line.fill.background()

    # Problem statement
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem Statement", "From isolated road sections to coordinated maintenance decisions at network scale")
    if problem_image and problem_image.exists():
        slide.shapes.add_picture(str(problem_image), Inches(0.65), Inches(1.2), width=Inches(12.0))
    add_textbox(slide, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.45), "Traditional pavement management predicts one section at a time; this dissertation asks how local deterioration, traffic, and dependencies can be linked to network-wide maintenance planning.", font_size=14, color=SLATE)

    # Results storyline
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Results Storyline", "What the results now support most clearly")
    add_bullet_box(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(11.2),
        Inches(4.8),
        [
            "Local history is strong, but graph information adds value when the graph semantics are well designed.",
            "On HPMS16, the strongest in-domain benchmark is the Stacked MLP ensemble (R² = 0.5356).",
            "A refined pavement-aware graph now reaches almost the same level on HPMS16 (R² = 0.5329).",
            "On MEPDG, the best absolute score is the refined R-GCN with materials + weight sweep (R² = 0.5582).",
            "Geographic transfer remains weak out-of-domain, which supports local recalibration rather than one universal model.",
        ],
        font_size=21,
    )

    # HPMS16 benchmark
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Core Benchmark on HPMS16")
    slide.shapes.add_picture(str(figs["hpms16"]), Inches(0.7), Inches(1.2), width=Inches(8.4))
    add_bullet_box(
        slide,
        Inches(9.3),
        Inches(1.45),
        Inches(3.2),
        Inches(4.6),
        [
            "RF local: 0.4489",
            "R-GCN full refined: 0.3708",
            "Refined R-GCN with materials + weights: 0.5329",
            "Stacked MLP ensemble: 0.5356",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(9.3), Inches(5.55), Inches(3.2), Inches(1.0), "Message: graph information matters most when it is either combined with strong local history, or encoded through a richer pavement similarity definition.", font_size=14, color=SLATE)

    # MEPDG refined
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Specialist Target: MEPDG Cracking")
    slide.shapes.add_picture(str(figs["mepdg"]), Inches(0.7), Inches(1.2), width=Inches(8.4))
    add_bullet_box(
        slide,
        Inches(9.25),
        Inches(1.45),
        Inches(3.3),
        Inches(4.6),
        [
            "RF local: 0.5233",
            "R-GCN baseline: 0.5424",
            "Stacked MLP ensemble: 0.5263",
            "Refined R-GCN: 0.5582",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(9.25), Inches(5.55), Inches(3.35), Inches(1.0), "Message: on MEPDG, the graph model itself is already stronger than RF, and the best score comes from refined graph semantics rather than ensembling.", font_size=14, color=SLATE)

    # OOD slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Robustness: Geographic Transfer")
    slide.shapes.add_picture(str(figs["ood"]), Inches(0.7), Inches(1.2), width=Inches(8.7))
    add_bullet_box(
        slide,
        Inches(9.5),
        Inches(1.45),
        Inches(3.0),
        Inches(4.8),
        [
            "Subset analysis on 5 representative held-out states",
            "Mean OOD R²: RF = 0.173",
            "Mean OOD R²: ensemble = 0.119",
            "Ensemble beats RF on only 2/5 states",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(9.5), Inches(5.5), Inches(3.0), Inches(1.0), "Message: the ensemble is a strong in-domain result, but it does not resolve state-to-state transfer. The signal remains region-specific.", font_size=14, color=SLATE)

    # Traffic / phase 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Bridge to Phase 2: Traffic Sensitivity")
    slide.shapes.add_picture(str(figs["traffic"]), Inches(0.65), Inches(1.2), width=Inches(9.0))
    add_bullet_box(
        slide,
        Inches(9.8),
        Inches(1.5),
        Inches(2.6),
        Inches(4.6),
        [
            "n events = 4,627",
            "n controls = 55,436",
            "Diff = -6,016",
            "95% CI = [-9,415 ; -2,536]",
            "Welch p = 0.0007",
        ],
        font_size=18,
    )
    add_textbox(slide, Inches(9.8), Inches(5.55), Inches(2.6), Inches(1.0), "Message: even with annual traffic-loading proxies, maintenance years leave measurable signatures that motivate the optimisation stage.", font_size=14, color=SLATE)

    # Final takeaways
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Takeaways")
    add_bullet_box(
        slide,
        Inches(0.9),
        Inches(1.45),
        Inches(11.3),
        Inches(4.8),
        [
            "Best comparative HPMS16 result: Stacked MLP ensemble, R² = 0.5356.",
            "Best absolute cracking result: refined R-GCN on MEPDG, R² = 0.5582.",
            "New result: refined pavement similarity also lifts HPMS16 strongly, to R² = 0.5329.",
            "Key limitation: geographic transfer remains weak, so local recalibration matters.",
            "Next step: translate these predictive and network signals into maintenance scheduling under disruption constraints.",
        ],
        font_size=22,
    )

    prs.save(PPTX_PATH)
    return PPTX_PATH


def main() -> None:
    setup_matplotlib()
    problem_image = copy_problem_statement_visual()
    figs = {
        "hpms16": build_hpms16_benchmark(),
        "mepdg": build_mepdg_benchmark(),
        "ood": build_ood_figure(),
        "traffic": build_traffic_figure(),
    }
    pptx_path = create_presentation(figs, problem_image)
    print(f"Saved {pptx_path}")
    for key, value in figs.items():
        print(f"{key}: {value}")
    if problem_image:
        print(f"problem_statement_image: {problem_image}")


if __name__ == "__main__":
    main()
